"""
Report Agent - 專用報告生成 Agent
使用 LangGraph 實作，支援 /report 指令觸發的報告生成流程
"""
import logging
import re
from datetime import datetime
from typing import Dict, Any, List, Optional, TypedDict, Annotated
from pathlib import Path

from langchain_core.messages import BaseMessage, HumanMessage, AIMessage
from langgraph.graph import StateGraph, END
from langgraph.graph.message import add_messages
try:
    from langchain.callbacks.tracers import LangChainTracer
except ImportError:
    # Fallback for different langchain versions
    try:
        from langchain_core.tracers import LangChainTracer
    except ImportError:
        # Create a dummy tracer if not available
        class LangChainTracer:
            def __init__(self, project_name=None):
                self.project_name = project_name

# 延遲導入，避免在直接執行時出現模組找不到的問題
settings = None
fmp_client = None
rag_service = None

# Shared tools from agent_graph (delayed import)
_agent_tool_fmp_quote = None
_agent_tool_fmp_profile = None
_agent_tool_fmp_news = None
_agent_tool_fmp_macro = None
_agent_tool_rag_query = None

# LLM analysis functions from agent_graph (delayed import)
_analyze_data_with_llm = None
_build_analysis_prompt = None

def _ensure_imports():
    """確保必要的模組已導入"""
    global settings, fmp_client, rag_service
    global _agent_tool_fmp_quote, _agent_tool_fmp_profile, _agent_tool_fmp_news, _agent_tool_fmp_macro, _agent_tool_rag_query
    global _analyze_data_with_llm, _build_analysis_prompt

    if settings is None:
        from app.settings import settings as _settings
        from app.services.fmp_client import fmp_client as _fmp_client
        from app.services.rag import rag_service as _rag_service
        settings = _settings
        fmp_client = _fmp_client
        rag_service = _rag_service

    # Import shared tools from agent_graph_simple to avoid duplication
    if _agent_tool_fmp_quote is None:
        try:
            from app.graphs.agent_graph_simple import (
                tool_fmp_quote as _tool_fmp_quote,
                tool_fmp_profile as _tool_fmp_profile,
                tool_fmp_news as _tool_fmp_news,
                tool_fmp_macro as _tool_fmp_macro,
                tool_rag_query as _tool_rag_query
            )
            _agent_tool_fmp_quote = _tool_fmp_quote
            _agent_tool_fmp_profile = _tool_fmp_profile
            _agent_tool_fmp_news = _tool_fmp_news
            _agent_tool_fmp_macro = _tool_fmp_macro
            _agent_tool_rag_query = _tool_rag_query

            # Create simple analysis functions since they don't exist in agent_graph
            def _simple_analyze_data_with_llm(data, prompt):
                """簡單的資料分析函數"""
                return f"分析結果：{str(data)[:200]}..."

            def _simple_build_analysis_prompt(symbols, data_type):
                """簡單的分析提示建構函數"""
                return f"請分析 {symbols} 的 {data_type} 資料"

            _analyze_data_with_llm = _simple_analyze_data_with_llm
            _build_analysis_prompt = _simple_build_analysis_prompt

        except ImportError as e:
            logger.warning(f"無法從 agent_graph_simple 導入工具：{e}")
            # 使用本地定義的工具
            _agent_tool_fmp_quote = tool_fmp_quote
            _agent_tool_fmp_profile = tool_fmp_profile
            _agent_tool_fmp_news = tool_fmp_news
            _agent_tool_fmp_macro = tool_fmp_macro
            _agent_tool_rag_query = tool_rag_query

            def _simple_analyze_data_with_llm(data, prompt):
                return f"分析結果：{str(data)[:200]}..."

            def _enhanced_build_analysis_prompt(context, template_id):
                """建構增強的分析提示，適用於各種資料源"""
                symbols = context.get("symbols", [])
                quotes = context.get("quotes", {})
                profiles = context.get("profiles", {})
                news = context.get("news", {})
                macro = context.get("macro", {})
                rag = context.get("rag", {})

                prompt_parts = []
                prompt_parts.append("請基於以下可用資料進行專業金融分析，並以 JSON 格式回應：")

                # 股票基本資訊
                if symbols:
                    prompt_parts.append(f"\n## 分析標的\n股票代號：{', '.join(symbols)}")

                # 股價資料
                if quotes.get("ok") and quotes.get("data"):
                    prompt_parts.append("\n## 股價資料")
                    for quote in quotes["data"]:
                        prompt_parts.append(f"- {quote.get('symbol', 'N/A')}: ${quote.get('price', 'N/A')}, 變動: {quote.get('change', 'N/A')} ({quote.get('changesPercentage', 'N/A')}%)")

                # 公司基本面
                if profiles.get("ok") and profiles.get("data"):
                    prompt_parts.append("\n## 公司基本面")
                    for profile in profiles["data"]:
                        prompt_parts.append(f"- {profile.get('symbol', 'N/A')}: {profile.get('companyName', 'N/A')}, 產業: {profile.get('industry', 'N/A')}")
                        if profile.get("description"):
                            prompt_parts.append(f"  描述: {profile['description'][:200]}...")

                # 新聞資料
                if news.get("ok") and news.get("data"):
                    prompt_parts.append("\n## 相關新聞")
                    for i, article in enumerate(news["data"][:3], 1):
                        prompt_parts.append(f"{i}. {article.get('title', 'N/A')}")
                        if article.get("text"):
                            prompt_parts.append(f"   摘要: {article['text'][:150]}...")

                # RAG 資料（如果可用）
                if rag.get("ok") and rag.get("data"):
                    prompt_parts.append("\n## 文件檢索資料")
                    prompt_parts.append(f"檢索結果: {str(rag['data'])[:300]}...")

                # 總經資料
                if macro.get("ok") and macro.get("data"):
                    prompt_parts.append("\n## 總體經濟資料")
                    prompt_parts.append(f"總經指標: {str(macro['data'])[:200]}...")

                # 分析要求
                prompt_parts.append("""
## 分析要求
請提供結構化的 JSON 分析，**必須**包含以下確切欄位名稱：
{
  "market_analysis": "基於股價變動和市場表現的分析（200字以內）",
  "fundamental_analysis": "基於公司基本面資料的分析（200字以內）",
  "news_impact": "基於相關新聞對股價可能影響的評估（150字以內）",
  "investment_recommendation": "基於綜合分析的投資建議（150字以內）",
  "risk_assessment": "投資風險評估和注意事項（150字以內）",
  "key_insights": ["關鍵洞察1", "關鍵洞察2", "關鍵洞察3"]
}

**重要要求：**
1. 回應必須是純 JSON 格式，不要包含其他文字
2. 欄位名稱必須完全一致：market_analysis, fundamental_analysis, news_impact, investment_recommendation, risk_assessment, key_insights
3. 僅基於提供的資料進行分析，不得杜撰
4. 如果某類資料不足，請在相應欄位中說明資料限制
5. 保持客觀和專業的分析態度
6. 所有文字使用繁體中文
7. key_insights 必須是字串陣列，包含 3-5 個要點""")

                return "\n".join(prompt_parts)

            _analyze_data_with_llm = _simple_analyze_data_with_llm
            _build_analysis_prompt = _enhanced_build_analysis_prompt

logger = logging.getLogger(__name__)


# ===== Report Agent Tools =====

async def tool_fmp_quote(symbols: List[str]) -> Dict[str, Any]:
    """FMP 股票報價查詢工具 - 使用 agent_graph 的共享實現"""
    _ensure_imports()
    logger.info(f"FMP 報價查詢: {symbols}")

    try:
        if not symbols:
            return {
                "ok": False,
                "source": "FMP",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "未提供股票代號",
                "error": "missing_symbols"
            }

        # 使用 agent_graph 的共享工具
        result = await _agent_tool_fmp_quote.ainvoke({"symbols": symbols})

        # 包裝結果以保持與現有格式的兼容性
        return {
            "ok": result.get("ok", False),
            "source": "FMP",
            "data": result.get("data"),
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢 {len(symbols)} 個股票報價",
            "error": result.get("reason") if not result.get("ok") else None
        }

    except Exception as e:
        logger.error(f"FMP 報價查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢失敗: {str(e)}",
            "error": "query_failed"
        }


async def tool_fmp_profile(symbols: List[str]) -> Dict[str, Any]:
    """FMP 公司基本面查詢工具 - 使用 agent_graph 的共享實現"""
    _ensure_imports()
    logger.info(f"FMP 基本面查詢: {symbols}")

    try:
        if not symbols:
            return {
                "ok": False,
                "source": "FMP",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "未提供股票代號",
                "error": "missing_symbols"
            }

        # 使用 agent_graph 的共享工具
        result = await _agent_tool_fmp_profile.ainvoke({"symbols": symbols})

        # 包裝結果以保持與現有格式的兼容性
        return {
            "ok": result.get("ok", False),
            "source": "FMP",
            "data": result.get("data"),
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢 {len(symbols)} 個公司基本面",
            "error": result.get("reason") if not result.get("ok") else None
        }

    except Exception as e:
        logger.error(f"FMP 基本面查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢失敗: {str(e)}",
            "error": "query_failed"
        }


async def tool_fmp_news(symbols: List[str], limit: int = 10) -> Dict[str, Any]:
    """FMP 新聞查詢工具 - 使用 agent_graph 的共享實現"""
    _ensure_imports()
    logger.info(f"FMP 新聞查詢: {symbols}, limit={limit}")

    try:
        # 使用 agent_graph 的共享工具
        # Note: agent_graph tool_fmp_news has different signature: (symbols=None, query=None, limit=10)
        result = await _agent_tool_fmp_news.ainvoke({
            "symbols": symbols if symbols else None,
            "query": None,
            "limit": limit
        })

        # 包裝結果以保持與現有格式的兼容性
        return {
            "ok": result.get("ok", False),
            "source": "FMP",
            "data": result.get("data"),
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢新聞，股票: {symbols if symbols else '一般市場'}, 限制: {limit}",
            "error": result.get("reason") if not result.get("ok") else None
        }

    except Exception as e:
        logger.error(f"FMP 新聞查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢失敗: {str(e)}",
            "error": "query_failed"
        }


async def tool_fmp_macro(indicators: List[str], limit: int = 6) -> Dict[str, Any]:
    """FMP 總經指標查詢工具 - 使用 agent_graph 的共享實現"""
    _ensure_imports()
    logger.info(f"FMP 總經查詢: {indicators}, limit={limit}")

    try:
        if not indicators:
            # 預設查詢主要指標
            indicators = ["GDP", "CPI", "UNEMPLOYMENT"]

        macro_data = {}

        # 使用 agent_graph 的共享工具
        # Note: agent_graph tool_fmp_macro has signature: (indicator: str, country: str = "US")
        # We need to call it for each indicator individually
        for indicator in indicators:
            try:
                result = await _agent_tool_fmp_macro.ainvoke({"indicator": indicator, "country": "US"})
                if result.get("ok"):
                    macro_data[f"{indicator}_US"] = result.get("data", [])
                else:
                    logger.warning(f"總經指標 {indicator} 查詢失敗: {result.get('reason')}")
                    macro_data[f"{indicator}_US"] = []
            except Exception as e:
                logger.error(f"總經指標 {indicator} 查詢異常: {str(e)}")
                macro_data[f"{indicator}_US"] = []

        return {
            "ok": True,
            "source": "FMP",
            "data": macro_data,
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢總經指標: {indicators}, 限制: {limit}",
            "error": None
        }

    except Exception as e:
        logger.error(f"FMP 總經查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢失敗: {str(e)}",
            "error": "query_failed"
        }


def tool_select_template(report_type: str, context: Dict[str, Any]) -> Dict[str, Any]:
    """模板選擇工具"""
    logger.info(f"模板選擇: {report_type}")

    try:
        # 模板映射
        template_mapping = {
            "stock": "stock.j2",
            "macro": "macro.j2",
            "news": "news.j2",
            "custom": "custom.j2"
        }

        template_id = template_mapping.get(report_type, "custom.j2")
        template_path = Path("templates/reports") / template_id

        # 檢查模板是否存在
        if not template_path.exists():
            return {
                "ok": False,
                "source": "TEMPLATE",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"模板檔案不存在: {template_path}",
                "error": "template_not_found"
            }

        return {
            "ok": True,
            "source": "TEMPLATE",
            "data": {
                "template_id": template_id,
                "template_path": str(template_path),
                "report_type": report_type
            },
            "timestamp": datetime.now().isoformat(),
            "logs": f"已選擇模板: {template_id}",
            "error": None
        }

    except Exception as e:
        logger.error(f"模板選擇失敗: {str(e)}")
        return {
            "ok": False,
            "source": "TEMPLATE",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"選擇失敗: {str(e)}",
            "error": "selection_failed"
        }


async def tool_rag_query(question: str, top_k: int = 5) -> Dict[str, Any]:
    """RAG 查詢工具 - 使用 agent_graph 的共享實現"""
    _ensure_imports()
    logger.info(f"RAG 查詢: {question}, top_k={top_k}")

    try:
        if not question.strip():
            return {
                "ok": False,
                "source": "RAG",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "查詢字串為空",
                "error": "empty_query"
            }

        # 使用 agent_graph 的共享工具
        result = await _agent_tool_rag_query.ainvoke({"question": question, "top_k": top_k})

        # 包裝結果以保持與現有格式的兼容性
        if result.get("ok"):
            return {
                "ok": True,
                "source": "RAG",
                "data": result.get("data"),
                "timestamp": datetime.now().isoformat(),
                "logs": f"RAG 查詢成功，返回 {top_k} 個結果",
                "error": None
            }
        else:
            return {
                "ok": False,
                "source": "RAG",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"RAG 查詢失敗: {result.get('reason', 'unknown')}",
                "error": result.get("reason", "query_failed")
            }

    except Exception as e:
        # 將 RAG 失敗從錯誤級別降級為警告級別
        logger.warning(f"RAG 查詢失敗，但不影響報告生成: {str(e)}")

        # 檢查是否是向量維度不匹配錯誤
        error_type = "query_failed"
        if "shapes" in str(e) and "not aligned" in str(e):
            error_type = "vector_dimension_mismatch"
            logger.warning("檢測到向量維度不匹配錯誤，這通常是 RAG 服務配置問題")
        elif "connection" in str(e).lower() or "timeout" in str(e).lower():
            error_type = "connection_failed"
            logger.warning("RAG 服務連接失敗，可能服務不可用")

        return {
            "ok": False,
            "source": "RAG",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"RAG 查詢不可用: {str(e)}",
            "error": error_type,
            "fallback_note": "RAG 資料源暫時不可用，報告將基於其他可用資料源生成"
        }


async def tool_build_report(template_id: str, context: Dict[str, Any], output_formats: List[str] = None) -> Dict[str, Any]:
    """報告建構工具"""
    _ensure_imports()
    logger.info(f"建構報告: {template_id}, 格式: {output_formats}")

    if output_formats is None:
        output_formats = ["markdown", "pdf"]

    try:
        from jinja2 import Environment, FileSystemLoader
        import markdown

        # 準備輸出目錄
        timestamp = context.get("timestamp", generate_timestamp())
        slug = context.get("slug", "REPORT")
        output_dir = Path(settings.output_dir) / "reports" / timestamp
        output_dir.mkdir(parents=True, exist_ok=True)

        # 載入 Jinja2 模板
        template_dir = Path("templates/reports")
        if not template_dir.exists():
            return {
                "ok": False,
                "source": "REPORT",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"模板目錄不存在: {template_dir}",
                "error": "template_dir_not_found"
            }

        env = Environment(loader=FileSystemLoader(str(template_dir)))

        try:
            template = env.get_template(template_id)
        except Exception as e:
            return {
                "ok": False,
                "source": "REPORT",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"載入模板失敗: {template_id}, 錯誤: {str(e)}",
                "error": "template_load_failed"
            }

        # 準備模板變數
        template_vars = {
            "generated_at": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "timestamp": timestamp,
            "slug": slug,
            **context
        }

        # 添加資料源可用性資訊
        data_source_status = {
            "rag_available": context.get("rag", {}).get("ok", False),
            "fmp_available": any([
                context.get("quotes", {}).get("ok", False),
                context.get("profiles", {}).get("ok", False),
                context.get("news", {}).get("ok", False),
                context.get("macro", {}).get("ok", False)
            ]),
            "data_source_notes": context.get("data_source_notes", [])
        }
        template_vars["data_source_status"] = data_source_status

        # 如果有資料源註記，添加到模板變數中
        if context.get("data_source_notes"):
            template_vars["has_data_limitations"] = True
            template_vars["data_limitations_note"] = "注意：" + "、".join(context["data_source_notes"])
        else:
            template_vars["has_data_limitations"] = False

        # 渲染模板
        try:
            rendered_content = template.render(**template_vars)
        except Exception as e:
            return {
                "ok": False,
                "source": "REPORT",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"模板渲染失敗: {str(e)}",
                "error": "template_render_failed"
            }

        output_files = []

        # 生成 Markdown 檔案
        if "markdown" in output_formats:
            md_file = output_dir / f"{slug}.md"
            try:
                with open(md_file, 'w', encoding='utf-8') as f:
                    f.write(rendered_content)

                output_files.append({
                    "format": "markdown",
                    "filename": f"{slug}.md",
                    "path": str(md_file),
                    "size": md_file.stat().st_size
                })
                logger.info(f"已生成 Markdown 檔案: {md_file}")
            except Exception as e:
                logger.error(f"生成 Markdown 檔案失敗: {str(e)}")

        # 生成 PDF 檔案
        if "pdf" in output_formats:
            try:
                import weasyprint

                pdf_file = output_dir / f"{slug}.pdf"

                # 轉換 Markdown 為 HTML
                html_content = markdown.markdown(rendered_content)

                # 載入 CSS 樣式
                css_path = Path(settings.pdf_default_css)
                css_content = ""
                if css_path.exists():
                    with open(css_path, 'r', encoding='utf-8') as f:
                        css_content = f.read()

                # 完整 HTML
                full_html = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{slug} 報告</title>
    <style>{css_content}</style>
</head>
<body>
    {html_content}
</body>
</html>
"""

                # 生成 PDF
                weasyprint.HTML(string=full_html).write_pdf(str(pdf_file))

                output_files.append({
                    "format": "pdf",
                    "filename": f"{slug}.pdf",
                    "path": str(pdf_file),
                    "size": pdf_file.stat().st_size
                })
                logger.info(f"已生成 PDF 檔案: {pdf_file}")

            except ImportError:
                logger.warning("WeasyPrint 未安裝，跳過 PDF 生成")
            except Exception as e:
                logger.error(f"生成 PDF 檔案失敗: {str(e)}")

        # 生成 PPTX 檔案（可選）
        if "pptx" in output_formats:
            try:
                from pptx import Presentation

                pptx_file = output_dir / f"{slug}.pptx"

                # 建立簡單的 PowerPoint
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = f"{slug} 報告"
                content.text = f"生成時間: {template_vars['generated_at']}\n\n{rendered_content[:500]}..."

                prs.save(str(pptx_file))

                output_files.append({
                    "format": "pptx",
                    "filename": f"{slug}.pptx",
                    "path": str(pptx_file),
                    "size": pptx_file.stat().st_size
                })
                logger.info(f"已生成 PPTX 檔案: {pptx_file}")

            except ImportError:
                logger.warning("python-pptx 未安裝，跳過 PPTX 生成")
            except Exception as e:
                logger.error(f"生成 PPTX 檔案失敗: {str(e)}")

        return {
            "ok": True,
            "source": "REPORT",
            "data": {
                "files": output_files,
                "output_dir": str(output_dir),
                "template_id": template_id
            },
            "timestamp": datetime.now().isoformat(),
            "logs": f"已生成 {len(output_files)} 個檔案",
            "error": None
        }

    except Exception as e:
        logger.error(f"報告建構失敗: {str(e)}")
        return {
            "ok": False,
            "source": "REPORT",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"建構失敗: {str(e)}",
            "error": "build_failed"
        }


class ReportAgentState(TypedDict):
    """Report Agent 狀態定義"""
    messages: Annotated[List[BaseMessage], add_messages]
    input_type: str            # "text"
    query: str                 # 原始 /report 指令
    report_type: Optional[str] # "stock" | "macro" | "news" | "custom"
    symbols: List[str]
    template_id: Optional[str]
    context: Dict[str, Any]    # 聚合資料容器
    output_files: List[Dict[str, Any]]
    timestamp: str             # 固定 ts；建立於流程開始
    session_id: Optional[str]


def parse_report_query(query: str) -> Dict[str, Any]:
    """解析 /report 指令"""
    logger.info(f"解析報告查詢: {query}")
    
    # 移除 /report 前綴
    content = query.replace("/report", "").strip()
    
    # 預設值
    result = {
        "report_type": "custom",
        "symbols": [],
        "indicators": [],
        "keywords": []
    }
    
    # 檢測報告類型
    if re.search(r'\b(stock|stocks|股票|個股|分析)\b', content, re.IGNORECASE):
        result["report_type"] = "stock"
    elif re.search(r'\b(macro|總經|經濟|gdp|cpi|inflation)\b', content, re.IGNORECASE):
        result["report_type"] = "macro"
    elif re.search(r'\b(news|新聞|消息)\b', content, re.IGNORECASE):
        result["report_type"] = "news"

    # 如果找到股票代號但沒有明確類型，預設為股票報告
    symbols = re.findall(r'\b[A-Z]{2,5}\b', content.upper())
    exclude_words = {'STOCK', 'STOCKS', 'MACRO', 'NEWS', 'CUSTOM', 'REPORT', 'GDP', 'CPI'}
    valid_symbols = [s for s in symbols if s not in exclude_words]

    if valid_symbols and result["report_type"] == "custom":
        result["report_type"] = "stock"
        logger.info(f"檢測到股票代號 {valid_symbols}，自動設定為股票報告")
    
    # 使用之前檢測的股票代號結果
    result["symbols"] = list(set(valid_symbols)) if 'valid_symbols' in locals() else []
    
    # 提取總經指標關鍵詞
    macro_keywords = re.findall(r'\b(GDP|CPI|INFLATION|UNEMPLOYMENT|INTEREST|RATE)\b', content.upper())
    result["indicators"] = list(set(macro_keywords))
    
    # 其他關鍵詞
    result["keywords"] = content.split()
    
    logger.info(f"解析結果: {result}")
    return result


def generate_timestamp() -> str:
    """生成時間戳記"""
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def generate_slug(symbols: List[str], report_type: str) -> str:
    """生成檔案名稱 slug"""
    if symbols:
        # 取前兩個 symbol
        slug_symbols = symbols[:2]
        return "_".join(slug_symbols)
    else:
        return "CUSTOM" if report_type == "custom" else report_type.upper()


class ReportAgent:
    """Report Agent 實作"""

    def __init__(self):
        _ensure_imports()
        self.llm = self._create_llm()
        self.graph = self._create_graph()
        self.tracer = LangChainTracer(project_name="agent")

    def _create_llm(self):
        """建立 LLM 實例"""
        try:
            from langchain_openai import ChatOpenAI

            if settings.openai_api_key:
                return ChatOpenAI(
                    model="gpt-4o-mini",
                    api_key=settings.openai_api_key,
                    temperature=0.3
                )
            elif settings.azure_openai_api_key and settings.azure_openai_endpoint:
                return ChatOpenAI(
                    model=settings.azure_openai_deployment,
                    api_key=settings.azure_openai_api_key,
                    azure_endpoint=settings.azure_openai_endpoint,
                    api_version=settings.azure_openai_api_version,
                    temperature=0.3
                )
            else:
                logger.warning("未設定 OpenAI 或 Azure OpenAI API 金鑰，LLM 增強功能將無法使用")
                return None
        except ImportError:
            logger.warning("未安裝 langchain_openai，LLM 增強功能將無法使用")
            return None

    async def _analyze_data_with_llm_local(self, context: Dict[str, Any], template_id: str) -> Dict[str, Any]:
        """使用本地 LLM 實例分析資料並增強 context"""
        try:
            if not self.llm:
                logger.warning("LLM 未設定，跳過分析")
                return context

            # 使用 agent_graph 的提示建構函數
            prompt = _build_analysis_prompt(context, template_id)

            # 建立訊息
            from langchain_core.messages import SystemMessage, HumanMessage
            messages = [
                SystemMessage(content="你是嚴謹的金融分析師，僅使用提供資料，產出結構化 JSON 洞察。不得杜撰。"),
                HumanMessage(content=prompt)
            ]

            # 調用 LLM（使用設定的參數）
            response = await self.llm.ainvoke(
                messages,
                temperature=settings.llm_analysis_temperature,
                max_tokens=settings.llm_analysis_max_tokens,
                timeout=settings.llm_analysis_timeout
            )

            # 處理回應（與 agent_graph 相同的邏輯）
            import json
            import re

            content = response.content.strip()
            logger.info(f"LLM 原始回應: {content[:200]}...")

            # 清理 JSON（移除 markdown 標記）
            content = re.sub(r'^```json\s*', '', content, flags=re.MULTILINE)
            content = re.sub(r'\s*```$', '', content, flags=re.MULTILINE)
            content = content.strip()

            logger.info(f"清理後的內容: {content[:200]}...")
            analysis_data = json.loads(content)
            logger.info(f"LLM 分析數據解析成功: {list(analysis_data.keys())}")

            # 智能提取分析欄位
            enhanced_context = {**context, "llm_analysis": analysis_data}

            # 定義需要的分析欄位
            analysis_fields = ["market_analysis", "fundamental_analysis", "news_impact",
                             "investment_recommendation", "risk_assessment", "key_insights"]

            # 首先檢查頂層是否有這些欄位
            for field in analysis_fields:
                if field in analysis_data:
                    enhanced_context[field] = analysis_data[field]

            # 如果頂層沒有，檢查 stock_analysis 嵌套結構
            if "stock_analysis" in analysis_data and not any(field in analysis_data for field in analysis_fields):
                stock_data = analysis_data["stock_analysis"]
                for field in analysis_fields:
                    if field in stock_data:
                        enhanced_context[field] = stock_data[field]

            # 如果仍然沒有找到分析欄位，使用基於可用資料的基本分析
            if not any(field in enhanced_context for field in analysis_fields):
                logger.warning("LLM 回應格式不符預期，生成基本分析")
                fallback_analysis = await self._provide_fallback_analysis(context)
                for field in analysis_fields:
                    if field in fallback_analysis:
                        enhanced_context[field] = fallback_analysis[field]
            logger.info(f"增強後的 context 包含: {list(enhanced_context.keys())}")
            return enhanced_context
        except json.JSONDecodeError as e:
            # JSON 解析失敗，將原始回應作為 AI 洞察
            logger.error(f"LLM 回應 JSON 解析失敗: {e}, 清理後內容: {content[:200] if 'content' in locals() else 'N/A'}...")
            return {
                **context,
                "llm_analysis": response.content,
                "ai_insights": response.content
            }
        except Exception as e:
            logger.error(f"LLM 分析過程異常：{e}")
            return context

    async def _provide_fallback_analysis(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """當 LLM 分析失敗時，提供基於可用資料的降級分析"""
        logger.info("提供降級分析...")

        symbols = context.get("symbols", [])
        quotes = context.get("quotes", {})
        profiles = context.get("profiles", {})
        news = context.get("news", {})

        # 基本市場分析
        market_analysis = "基於可用資料的市場分析："
        if quotes.get("ok") and quotes.get("data"):
            positive_changes = sum(1 for q in quotes["data"] if q.get("change", 0) > 0)
            total_stocks = len(quotes["data"])

            # 詳細分析每支股票
            stock_details = []
            for quote in quotes["data"]:
                symbol = quote.get("symbol", "N/A")
                price = quote.get("price", 0)
                change = quote.get("change", 0)
                change_pct = quote.get("changesPercentage", 0)

                trend = "上漲" if change > 0 else "下跌" if change < 0 else "持平"
                stock_details.append(f"{symbol} 目前價格 ${price:.2f}，{trend} {abs(change):.2f} ({change_pct:.2f}%)")

            market_analysis += f" {'; '.join(stock_details)}。"

            if positive_changes > total_stocks / 2:
                market_analysis += "整體而言，多數股票呈現上漲趨勢，市場情緒相對樂觀。"
            else:
                market_analysis += "整體而言，多數股票呈現下跌趨勢，市場情緒偏向謹慎。"
        else:
            market_analysis += "股價資料不可用，無法進行市場趨勢分析。"

        # 基本面分析
        fundamental_analysis = "基本面分析："
        if profiles.get("ok") and profiles.get("data"):
            company_details = []
            for profile in profiles["data"]:
                symbol = profile.get("symbol", "N/A")
                company_name = profile.get("companyName", "N/A")
                industry = profile.get("industry", "未知")
                market_cap = profile.get("mktCap", 0)
                employees = profile.get("fullTimeEmployees", 0)

                # 格式化市值
                if market_cap > 1e12:
                    market_cap_str = f"{market_cap/1e12:.1f}兆美元"
                elif market_cap > 1e9:
                    market_cap_str = f"{market_cap/1e9:.1f}億美元"
                else:
                    market_cap_str = f"{market_cap:.0f}美元"

                company_details.append(f"{symbol} ({company_name}) 屬於{industry}產業，市值約{market_cap_str}")
                if employees > 0:
                    company_details[-1] += f"，員工數約{employees:,}人"

            fundamental_analysis += f" {'; '.join(company_details)}。"

            # 產業分析
            industries = [p.get("industry", "未知") for p in profiles["data"]]
            unique_industries = list(set(industries))
            if len(unique_industries) > 1:
                fundamental_analysis += f"投資組合涵蓋{len(unique_industries)}個不同產業，具有一定的分散效果。"
            else:
                fundamental_analysis += f"投資集中在{unique_industries[0]}產業。"
        else:
            fundamental_analysis += "公司基本面資料不可用，建議查詢公司財報和基本資訊。"

        # 新聞影響評估
        news_impact = "新聞影響評估："
        if news.get("ok") and news.get("data"):
            news_count = len(news["data"])
            recent_news = news["data"][:3]  # 取前3則新聞

            news_titles = []
            for article in recent_news:
                title = article.get("title", "無標題")
                date = article.get("publishedDate", "")
                if date:
                    # 簡化日期格式
                    try:
                        from datetime import datetime
                        date_obj = datetime.fromisoformat(date.replace('Z', '+00:00'))
                        date_str = date_obj.strftime("%m/%d")
                    except:
                        date_str = date[:10] if len(date) >= 10 else date
                    news_titles.append(f"({date_str}) {title}")
                else:
                    news_titles.append(title)

            news_impact += f"共發現 {news_count} 則相關新聞。"
            if news_titles:
                news_impact += f"近期重要新聞包括：{'; '.join(news_titles[:2])}等。"
            news_impact += "建議關注新聞動態對股價的潛在影響。"
        else:
            news_impact += "新聞資料不可用，建議手動查詢相關新聞以評估媒體影響。"

        # 投資建議
        investment_recommendation = "投資建議："
        if quotes.get("ok") and quotes.get("data"):
            positive_stocks = [q for q in quotes["data"] if q.get("change", 0) > 0]
            if positive_stocks:
                investment_recommendation += f"目前有 {len(positive_stocks)} 支股票呈現上漲趨勢，"
                if len(positive_stocks) == len(quotes["data"]):
                    investment_recommendation += "整體表現良好，可考慮適度投資。"
                else:
                    investment_recommendation += "建議重點關注表現較佳的標的。"
            else:
                investment_recommendation += "目前股票表現偏弱，建議謹慎觀望或考慮分批進場。"
        else:
            investment_recommendation += "由於缺乏即時股價資料，"

        investment_recommendation += "建議投資者結合基本面分析、技術分析和專業分析師意見進行綜合決策。"

        # 風險評估
        risk_assessment = "風險評估："
        risk_factors = []

        if quotes.get("ok") and quotes.get("data"):
            volatile_stocks = [q for q in quotes["data"] if abs(q.get("changesPercentage", 0)) > 3]
            if volatile_stocks:
                risk_factors.append(f"有 {len(volatile_stocks)} 支股票波動較大（超過3%）")

        if profiles.get("ok") and profiles.get("data"):
            industries = [p.get("industry", "") for p in profiles["data"]]
            unique_industries = list(set(industries))
            if len(unique_industries) == 1:
                risk_factors.append("投資集中在單一產業，缺乏分散效果")

        if not news.get("ok"):
            risk_factors.append("缺乏新聞資訊，可能錯過重要市場動態")

        if risk_factors:
            risk_assessment += f"主要風險包括：{'; '.join(risk_factors)}。"
        else:
            risk_assessment += "基於可用資料，"

        risk_assessment += "建議分散投資、設定停損點，並密切關注市場變化。"

        # 關鍵洞察
        key_insights = []

        if quotes.get("ok") and quotes.get("data"):
            avg_change = sum(q.get("changesPercentage", 0) for q in quotes["data"]) / len(quotes["data"])
            if avg_change > 1:
                key_insights.append("整體股票表現良好，平均漲幅超過1%")
            elif avg_change < -1:
                key_insights.append("整體股票表現疲弱，平均跌幅超過1%")
            else:
                key_insights.append("股票表現相對穩定，波動較小")

        if profiles.get("ok") and profiles.get("data"):
            large_caps = [p for p in profiles["data"] if p.get("mktCap", 0) > 1e11]  # 超過1000億美元
            if large_caps:
                key_insights.append(f"投資組合包含 {len(large_caps)} 支大型股，相對穩健")

        if news.get("ok") and news.get("data"):
            key_insights.append(f"近期有 {len(news['data'])} 則相關新聞，需關注媒體影響")

        # 確保至少有基本洞察
        if not key_insights:
            key_insights = [
                "基於當前可用資料進行分析",
                "建議結合多方資訊進行投資決策",
                "市場變化快速，請持續關注最新動態"
            ]

        # 添加分析結果到上下文
        fallback_context = {
            **context,
            "market_analysis": market_analysis,
            "fundamental_analysis": fundamental_analysis,
            "news_impact": news_impact,
            "investment_recommendation": investment_recommendation,
            "risk_assessment": risk_assessment,
            "key_insights": key_insights,
            "llm_analysis": {
                "market_analysis": market_analysis,
                "fundamental_analysis": fundamental_analysis,
                "news_impact": news_impact,
                "investment_recommendation": investment_recommendation,
                "risk_assessment": risk_assessment,
                "key_insights": key_insights,
                "analysis_type": "fallback_analysis"
            }
        }

        logger.info("降級分析完成")
        return fallback_context


    def _create_graph(self) -> StateGraph:
        """建立 LangGraph"""
        workflow = StateGraph(ReportAgentState)
        
        # 加入節點
        workflow.add_node("parse_query", self.parse_query_node)
        workflow.add_node("collect_data", self.collect_data_node)
        workflow.add_node("select_template", self.select_template_node)
        workflow.add_node("build_report", self.build_report_node)
        workflow.add_node("finalize", self.finalize_node)
        
        # 設定流程
        workflow.set_entry_point("parse_query")
        workflow.add_edge("parse_query", "collect_data")
        workflow.add_edge("collect_data", "select_template")
        workflow.add_edge("select_template", "build_report")
        workflow.add_edge("build_report", "finalize")
        workflow.add_edge("finalize", END)
        
        return workflow.compile()
    
    def parse_query_node(self, state: ReportAgentState) -> ReportAgentState:
        """解析查詢節點"""
        logger.info("執行查詢解析")
        
        query = state["query"]
        parsed = parse_report_query(query)
        
        # 生成固定時間戳記
        timestamp = generate_timestamp()
        
        # 更新狀態
        state.update({
            "report_type": parsed["report_type"],
            "symbols": parsed["symbols"],
            "timestamp": timestamp,
            "context": {
                "parsed_query": parsed,
                "slug": generate_slug(parsed["symbols"], parsed["report_type"])
            }
        })
        
        # 加入解析訊息
        state["messages"].append(
            AIMessage(content=f"已解析報告請求：類型={parsed['report_type']}, 股票代號={parsed['symbols']}")
        )
        
        return state
    
    async def collect_data_node(self, state: ReportAgentState) -> ReportAgentState:
        """資料收集節點 - 強制執行工具"""
        logger.info("執行資料收集")

        report_type = state["report_type"]
        symbols = state["symbols"]
        context = state["context"]

        # 檢查是否啟用工具執行
        if not settings.execute_tools:
            logger.warning("工具執行已停用")
            context["error"] = {
                "ok": False,
                "source": "SYSTEM",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "工具執行已停用",
                "error": "execute_tools_disabled"
            }
            state["context"] = context
            return state

        try:
            # 根據報告類型收集資料
            if report_type == "stock" and symbols:
                # 收集股票資料
                logger.info(f"收集股票資料: {symbols}")

                # 獲取報價
                quotes_result = await tool_fmp_quote(symbols)
                context["quotes"] = quotes_result

                # 獲取公司基本面
                profiles_result = await tool_fmp_profile(symbols)
                context["profiles"] = profiles_result

                # 獲取新聞
                news_result = await tool_fmp_news(symbols, limit=10)
                context["news"] = news_result

            elif report_type == "macro":
                # 收集總經資料
                logger.info("收集總經資料")
                indicators = ["GDP", "CPI", "UNEMPLOYMENT"]
                macro_result = await tool_fmp_macro(indicators, limit=6)
                context["macro"] = macro_result

                # 也收集相關新聞
                news_result = await tool_fmp_news([], limit=5)
                context["news"] = news_result

            elif report_type == "news":
                # 收集新聞資料
                logger.info(f"收集新聞資料: {symbols}")
                news_result = await tool_fmp_news(symbols if symbols else [], limit=15)
                context["news"] = news_result

                # 如果有股票代號，也獲取報價
                if symbols:
                    quotes_result = await tool_fmp_quote(symbols)
                    context["quotes"] = quotes_result

            else:
                # 自訂報告 - 嘗試使用 RAG，但不讓失敗阻止報告生成
                logger.info("收集 RAG 資料")
                try:
                    rag_result = await tool_rag_query(question=state["query"], top_k=5)
                    context["rag"] = rag_result

                    if not rag_result.get("ok"):
                        logger.warning(f"RAG 查詢失敗，將依賴其他資料源: {rag_result.get('error', 'unknown')}")
                        # 添加降級註記到上下文
                        context["data_source_notes"] = context.get("data_source_notes", [])
                        context["data_source_notes"].append("RAG 資料源暫時不可用")

                except Exception as e:
                    logger.warning(f"RAG 查詢異常，跳過此資料源: {str(e)}")
                    context["rag"] = {
                        "ok": False,
                        "source": "RAG",
                        "data": None,
                        "error": "service_unavailable",
                        "fallback_note": "RAG 服務暫時不可用"
                    }
                    context["data_source_notes"] = context.get("data_source_notes", [])
                    context["data_source_notes"].append("RAG 資料源暫時不可用")

                # 如果有股票代號，也獲取相關資料
                if symbols:
                    quotes_result = await tool_fmp_quote(symbols)
                    context["quotes"] = quotes_result

                    news_result = await tool_fmp_news(symbols, limit=3)
                    context["news"] = news_result

            # Fallback 機制：如果沒有收集到任何資料，嘗試其他工具
            has_data = any([
                context.get("quotes", {}).get("ok"),
                context.get("profiles", {}).get("ok"),
                context.get("news", {}).get("ok"),
                context.get("macro", {}).get("ok"),
                context.get("rag", {}).get("ok")
            ])

            if not has_data:
                logger.warning("未收集到任何資料，啟動 fallback 機制")

                # 如果有股票代號，嘗試獲取基本資料
                if symbols:
                    logger.info(f"Fallback: 嘗試獲取 {symbols} 的基本資料")
                    quotes_result = await tool_fmp_quote(symbols)
                    context["quotes"] = quotes_result

                    if not quotes_result.get("ok"):
                        # 如果股價也失敗，嘗試新聞
                        news_result = await tool_fmp_news(symbols, limit=5)
                        context["news"] = news_result

                # 否則嘗試 RAG 查詢（但不讓失敗阻止報告生成）
                else:
                    logger.info("Fallback: 嘗試 RAG 查詢")
                    try:
                        rag_result = await tool_rag_query(question=state["query"], top_k=3)
                        context["rag"] = rag_result

                        if not rag_result.get("ok"):
                            logger.warning("Fallback RAG 查詢也失敗，將生成基礎報告")
                            context["data_source_notes"] = context.get("data_source_notes", [])
                            context["data_source_notes"].append("所有資料源暫時不可用，生成基礎報告")

                    except Exception as e:
                        logger.warning(f"Fallback RAG 查詢異常: {str(e)}")
                        context["rag"] = {
                            "ok": False,
                            "source": "RAG",
                            "data": None,
                            "error": "fallback_failed",
                            "fallback_note": "RAG 服務在 fallback 階段也不可用"
                        }
                        context["data_source_notes"] = context.get("data_source_notes", [])
                        context["data_source_notes"].append("所有資料源暫時不可用，生成基礎報告")

        except Exception as e:
            # 將資料收集失敗從錯誤降級為警告，不阻止報告生成
            logger.warning(f"資料收集過程中發生異常，但將繼續生成報告: {str(e)}")
            context["error"] = {
                "ok": False,
                "source": "SYSTEM",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"資料收集異常: {str(e)}",
                "error": "data_collection_failed"
            }
            # 添加系統錯誤註記
            context["data_source_notes"] = context.get("data_source_notes", [])
            context["data_source_notes"].append("資料收集過程中發生異常，報告可能不完整")

        # 更新上下文
        context["symbols"] = symbols
        state["context"] = context

        # 加入收集訊息
        data_summary = []
        for key, value in context.items():
            if isinstance(value, dict) and "ok" in value:
                status = "成功" if value["ok"] else "失敗"
                data_summary.append(f"{key}: {status}")

        state["messages"].append(
            AIMessage(content=f"資料收集完成：{', '.join(data_summary)}")
        )

        return state
    
    def select_template_node(self, state: ReportAgentState) -> ReportAgentState:
        """模板選擇節點 - 支援 LLM 增強模板"""
        logger.info("執行模板選擇")
        _ensure_imports()

        report_type = state["report_type"]

        # 基礎模板映射
        template_mapping = {
            "stock": "stock.j2",
            "macro": "macro.j2",
            "news": "news.j2",
            "custom": "custom.j2"
        }

        base_template_id = template_mapping.get(report_type, "custom.j2")

        # 檢查是否啟用 LLM 增強並選擇對應模板
        if settings.llm_report_enhancement and report_type == "stock":
            # 對於股票報告，如果啟用 LLM 增強，使用增強模板
            template_id = "stock_llm_enhanced.j2"
            logger.info(f"LLM 增強已啟用，使用增強模板: {template_id}")
        elif report_type == "stock":
            # 即使 LLM 增強未啟用，股票報告也使用增強模板（提供更好的結構）
            template_id = "stock_llm_enhanced.j2"
            logger.info(f"股票報告使用增強模板（結構更完整）: {template_id}")
        else:
            template_id = base_template_id
            logger.info(f"使用基礎模板: {template_id}")

        state["template_id"] = template_id

        # 加入模板選擇訊息
        enhancement_status = "LLM增強=已啟用" if settings.llm_report_enhancement and report_type == "stock" else "LLM增強=未啟用"
        state["messages"].append(
            AIMessage(content=f"已選擇模板：{template_id} ({enhancement_status})")
        )

        return state
    
    async def build_report_node(self, state: ReportAgentState) -> ReportAgentState:
        """報告建構節點 - 支援 LLM 增強分析"""
        logger.info("執行報告建構")
        _ensure_imports()

        template_id = state["template_id"]
        context = state["context"]
        timestamp = state["timestamp"]

        # 準備基礎模板變數
        template_context = {
            **context,
            "timestamp": timestamp,
            "generated_at": datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }

        # LLM 增強分析（對於股票報告始終嘗試）
        if "llm_enhanced" in template_id:
            logger.info("開始 LLM 增強分析...")

            # 檢查是否有 LLM 可用
            if settings.llm_report_enhancement and self.llm:
                try:
                    # 使用本地 LLM 分析功能
                    enhanced_context = await self._analyze_data_with_llm_local(template_context, template_id)
                    template_context = enhanced_context
                    logger.info("LLM 增強分析完成")

                    # 加入 LLM 增強狀態訊息
                    state["messages"].append(
                        AIMessage(content="LLM 增強分析完成，正在生成增強報告...")
                    )
                except Exception as e:
                    logger.warning(f"LLM 增強分析失敗，使用降級分析: {str(e)}")

                    # 提供基本的分析結構
                    template_context = await self._provide_fallback_analysis(template_context)

                    state["messages"].append(
                        AIMessage(content="LLM 增強分析失敗，使用降級分析生成報告...")
                    )
            else:
                logger.info("LLM 不可用，使用降級分析")
                # 即使沒有 LLM，也提供基本分析結構
                template_context = await self._provide_fallback_analysis(template_context)

                state["messages"].append(
                    AIMessage(content="使用降級分析生成報告...")
                )
        else:
            logger.info("非增強模板，跳過 LLM 分析")

        # 使用報告建構工具
        build_result = await tool_build_report(
            template_id=template_id,
            context=template_context,
            output_formats=["markdown", "pdf"]
        )

        if build_result["ok"]:
            state["output_files"] = build_result["data"]["files"]

            # 加入成功訊息
            file_count = len(build_result["data"]["files"])
            state["messages"].append(
                AIMessage(content=f"報告建構成功，生成 {file_count} 個檔案")
            )
        else:
            # 建構失敗，記錄錯誤
            state["output_files"] = []
            error_msg = build_result.get("error", "unknown_error")

            state["messages"].append(
                AIMessage(content=f"報告建構失敗：{error_msg}")
            )

            logger.error(f"報告建構失敗: {build_result}")

        return state
    
    def finalize_node(self, state: ReportAgentState) -> ReportAgentState:
        """完成節點"""
        logger.info("執行報告完成")
        
        output_files = state["output_files"]
        file_count = len(output_files)
        
        # 加入完成訊息
        state["messages"].append(
            AIMessage(content=f"報告生成完成！共生成 {file_count} 個檔案。")
        )
        
        return state
    
    def _build_basic_report(self, state: ReportAgentState) -> str:
        """建構基本報告內容"""
        report_type = state["report_type"]
        symbols = state["symbols"]
        timestamp = state["timestamp"]
        
        # 基本報告模板
        content = f"""# {report_type.upper()} 報告

## 基本資訊
- 報告類型：{report_type}
- 生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
- 股票代號：{', '.join(symbols) if symbols else '無'}

## 內容摘要
此為 {report_type} 類型的報告。

## 資料來源
- 系統：Agent-Only LangGraph Service
- 時間戳記：{timestamp}

## 注意事項
本報告由 AI 系統自動生成，僅供參考。
"""
        return content
    
    async def run(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """執行 Report Agent"""
        try:
            # 準備初始狀態
            initial_state = ReportAgentState(
                messages=[HumanMessage(content=input_data["query"])],
                input_type=input_data["input_type"],
                query=input_data["query"],
                report_type=None,
                symbols=[],
                template_id=None,
                context={},
                output_files=[],
                timestamp="",
                session_id=input_data.get("session_id")
            )
            # 執行圖
            config = {"recursion_limit": 10}
            result = await self.graph.ainvoke(initial_state, config=config)
            
            # 建構回應
            return {
                "ok": True,
                "response": result["messages"][-1].content,
                "input_type": result["input_type"],
                "query": result["query"],
                "output_files": result["output_files"],
                "session_id": result["session_id"],
                "timestamp": datetime.now().isoformat(),
                "trace_id": input_data.get("trace_id", "")
            }
            
        except Exception as e:
            logger.error(f"Report Agent 執行失敗: {str(e)}")
            return {
                "ok": False,
                "response": f"報告生成失敗：{str(e)}",
                "input_type": input_data["input_type"],
                "query": input_data["query"],
                "output_files": [],
                "session_id": input_data.get("session_id"),
                "timestamp": datetime.now().isoformat(),
                "trace_id": input_data.get("trace_id", ""),
                "error": str(e)
            }


# 全域 Report Agent 實例（延遲初始化）
report_agent = None


# ===== Direct Execution Support =====

async def run_default_report_generation():
    """執行預設報告生成（無需命令列參數）"""
    from datetime import datetime
    global report_agent

    # 初始化 Report Agent（如果尚未初始化）
    if report_agent is None:
        report_agent = ReportAgent()

    # 硬編碼的預設值
    DEFAULT_QUERY = "/report stock AAPL TSLA NVDA"
    DEFAULT_OUTPUT_DIR = "./outputs"
    DEFAULT_FORMATS = ["markdown", "pdf"]

    print("🤖 Report Agent - 預設報告生成")
    print(f"📝 查詢：{DEFAULT_QUERY}")
    print(f"📁 輸出目錄：{DEFAULT_OUTPUT_DIR}")
    print(f"📄 輸出格式：{', '.join(DEFAULT_FORMATS)}")
    print()

    try:
        # 確保輸出目錄存在
        from pathlib import Path
        output_path = Path(DEFAULT_OUTPUT_DIR)
        output_path.mkdir(parents=True, exist_ok=True)

        # 準備輸入資料
        input_data = {
            "input_type": "text",
            "query": DEFAULT_QUERY.strip(),
            "session_id": f"default-{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            "trace_id": f"report-default-{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            "options": {
                "output_dir": DEFAULT_OUTPUT_DIR,
                "formats": DEFAULT_FORMATS
            }
        }

        print("🚀 開始執行報告生成...")

        # 執行報告生成
        result = await report_agent.run(input_data)

        # 處理結果
        if result.get("ok"):
            print("✅ 報告生成成功！")
            print()

            # 顯示輸出檔案
            output_files = result.get("output_files", [])
            if output_files:
                print("📄 生成的檔案：")
                for file_info in output_files:
                    file_path = file_info.get("path", "未知路徑")
                    file_format = file_info.get("format", "未知格式")
                    print(f"  • {file_format.upper()}: {file_path}")

            # 顯示回應內容
            if result.get("response"):
                print()
                print("📋 報告摘要：")
                print(result["response"])

            return 0

        else:
            print("❌ 報告生成失敗")
            error_msg = result.get("error", "未知錯誤")
            print(f"錯誤：{error_msg}")

            if result.get("response"):
                print()
                print("詳細資訊：")
                print(result["response"])

            return 1

    except KeyboardInterrupt:
        print("\n⏹️ 使用者中斷執行")
        return 130

    except Exception as e:
        print(f"❌ 執行失敗：{str(e)}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    """直接執行報告生成的主入口"""
    import asyncio
    import sys
    import os
    from pathlib import Path

    # 取得當前檔案的路徑
    current_file = Path(__file__).resolve()

    # 找到專案根目錄（包含 app 目錄的目錄）
    project_root = current_file.parent.parent.parent

    # 加入專案根目錄到 Python 路徑
    sys.path.insert(0, str(project_root))

    # 設定工作目錄為專案根目錄
    os.chdir(project_root)

    # 執行預設報告生成
    exit_code = asyncio.run(run_default_report_generation())
    sys.exit(exit_code)
