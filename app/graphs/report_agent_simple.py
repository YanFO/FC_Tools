"""
簡化報告代理 - 使用 LLM 自主決策系統
重新實現 report_agent.py 的完整功能，採用線性處理流程和智能查詢解析
所有註解、文字內容、日誌訊息使用繁體中文 (zh-TW)
"""
import logging
import re
import json
from datetime import datetime
from typing import Dict, Any, List, Optional, TypedDict, Annotated
from pathlib import Path

from langchain_core.messages import BaseMessage, HumanMessage, AIMessage, SystemMessage
from langchain_core.tools import tool

logger = logging.getLogger(__name__)

# 延遲導入的全域變數
_imports_loaded = False
_agent_tool_fmp_quote = None
_agent_tool_fmp_profile = None
_agent_tool_fmp_news = None
_agent_tool_fmp_macro = None
_agent_tool_rag_query = None
_analyze_data_with_llm = None
_build_analysis_prompt = None


def _ensure_imports():
    """確保所有必要的導入已載入"""
    global _imports_loaded, _agent_tool_fmp_quote, _agent_tool_fmp_profile
    global _agent_tool_fmp_news, _agent_tool_fmp_macro, _agent_tool_rag_query
    global _analyze_data_with_llm, _build_analysis_prompt

    if _imports_loaded:
        return

    try:
        # 導入設定
        from app.settings import settings

        # 嘗試從 agent_graph_simple 導入工具
        try:
            from app.graphs.agent_graph_simple import (
                tool_fmp_quote as agent_tool_fmp_quote,
                tool_fmp_profile as agent_tool_fmp_profile,
                tool_fmp_news as agent_tool_fmp_news,
                tool_fmp_macro as agent_tool_fmp_macro,
                tool_rag_query as agent_tool_rag_query
            )

            _agent_tool_fmp_quote = agent_tool_fmp_quote
            _agent_tool_fmp_profile = agent_tool_fmp_profile
            _agent_tool_fmp_news = agent_tool_fmp_news
            _agent_tool_fmp_macro = agent_tool_fmp_macro
            _agent_tool_rag_query = agent_tool_rag_query

            logger.info("成功從 agent_graph_simple 導入工具")

        except ImportError as e:
            logger.warning(f"無法從 agent_graph_simple 導入工具，使用本地實現: {e}")
            # 使用本地工具實現
            _agent_tool_fmp_quote = tool_fmp_quote
            _agent_tool_fmp_profile = tool_fmp_profile
            _agent_tool_fmp_news = tool_fmp_news
            _agent_tool_fmp_macro = tool_fmp_macro
            _agent_tool_rag_query = tool_rag_query

        # 設定 LLM 分析函數
        def _simple_analyze_data_with_llm(data, prompt):
            """簡單的 LLM 分析函數"""
            return f"基於可用資料的分析結果：{str(data)[:200]}..."

        def _simple_build_analysis_prompt(symbols, data_type):
            """簡單的分析提示建構函數"""
            return f"請分析 {symbols} 的 {data_type} 資料"

        _analyze_data_with_llm = _simple_analyze_data_with_llm
        _build_analysis_prompt = _simple_build_analysis_prompt

        _imports_loaded = True
        logger.info("導入初始化完成")

    except Exception as e:
        logger.error(f"導入初始化失敗: {str(e)}")
        raise


# ===== 本地工具實現 (備用) =====

@tool
async def tool_fmp_quote(symbols: List[str]) -> Dict[str, Any]:
    """FMP 股票報價查詢工具"""
    logger.info(f"FMP 報價查詢: {symbols}")

    try:
        from app.services.fmp_client import fmp_client

        if not symbols:
            return {
                "ok": False,
                "source": "FMP",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "未提供股票代號",
                "error": "missing_symbols"
            }

        result = await fmp_client.get_quote(symbols)

        return {
            "ok": result.get("ok", False),
            "source": "FMP",
            "data": result.get("data"),
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢 {len(symbols)} 個股票報價",
            "error": result.get("reason") if not result.get("ok") else None
        }

    except Exception as e:
        logger.error(f"FMP 報價查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"FMP 報價查詢失敗: {str(e)}",
            "error": "query_failed"
        }


@tool
async def tool_fmp_profile(symbols: List[str]) -> Dict[str, Any]:
    """FMP 公司基本面查詢工具"""
    logger.info(f"FMP 基本面查詢: {symbols}")

    try:
        from app.services.fmp_client import fmp_client

        if not symbols:
            return {
                "ok": False,
                "source": "FMP",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "未提供股票代號",
                "error": "missing_symbols"
            }

        result = await fmp_client.get_profile(symbols)

        return {
            "ok": result.get("ok", False),
            "source": "FMP",
            "data": result.get("data"),
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢 {len(symbols)} 個公司基本面",
            "error": result.get("reason") if not result.get("ok") else None
        }

    except Exception as e:
        logger.error(f"FMP 基本面查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"FMP 基本面查詢失敗: {str(e)}",
            "error": "query_failed"
        }


@tool
async def tool_fmp_news(symbols: Optional[List[str]] = None, limit: int = 10) -> Dict[str, Any]:
    """FMP 新聞查詢工具"""
    logger.info(f"FMP 新聞查詢: symbols={symbols}, limit={limit}")

    try:
        from app.services.fmp_client import fmp_client

        result = await fmp_client.get_news(symbols=symbols, limit=limit)

        return {
            "ok": result.get("ok", False),
            "source": "FMP",
            "data": result.get("data"),
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢 {limit} 條新聞",
            "error": result.get("reason") if not result.get("ok") else None
        }

    except Exception as e:
        logger.error(f"FMP 新聞查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"FMP 新聞查詢失敗: {str(e)}",
            "error": "query_failed"
        }


@tool
async def tool_fmp_macro(indicators: List[str], limit: int = 6) -> Dict[str, Any]:
    """FMP 總經指標查詢工具"""
    logger.info(f"FMP 總經查詢: indicators={indicators}, limit={limit}")

    try:
        from app.services.fmp_client import fmp_client

        if not indicators:
            indicators = ["GDP", "CPI", "UNEMPLOYMENT"]

        macro_data = {}

        # 逐個查詢指標
        for indicator in indicators:
            try:
                result = await fmp_client.get_macro_data(indicator, "US")
                if result.get("ok"):
                    macro_data[f"{indicator}_US"] = result.get("data", [])
                else:
                    logger.warning(f"總經指標 {indicator} 查詢失敗: {result.get('reason')}")
                    macro_data[f"{indicator}_US"] = []
            except Exception as e:
                logger.error(f"總經指標 {indicator} 查詢異常: {str(e)}")
                macro_data[f"{indicator}_US"] = []

        return {
            "ok": True,
            "source": "FMP",
            "data": macro_data,
            "timestamp": datetime.now().isoformat(),
            "logs": f"查詢 {len(indicators)} 個總經指標",
            "error": None
        }

    except Exception as e:
        logger.error(f"FMP 總經查詢失敗: {str(e)}")
        return {
            "ok": False,
            "source": "FMP",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"FMP 總經查詢失敗: {str(e)}",
            "error": "query_failed"
        }


@tool
async def tool_rag_query(question: str, top_k: int = 5) -> Dict[str, Any]:
    """RAG 查詢工具"""
    logger.info(f"RAG 查詢: {question}")

    try:
        from app.services.rag import rag_service

        query_result = await rag_service.query_documents(question, top_k)

        if query_result["ok"] and query_result["data"]["relevant_chunks"]:
            answer_result = await rag_service.answer_question(question, query_result["data"]["relevant_chunks"])
            return {
                "ok": answer_result.get("ok", False),
                "source": "RAG",
                "data": answer_result.get("data"),
                "timestamp": datetime.now().isoformat(),
                "logs": f"RAG 查詢完成，找到 {len(query_result['data']['relevant_chunks'])} 個相關片段",
                "error": answer_result.get("reason") if not answer_result.get("ok") else None
            }
        else:
            return {
                "ok": False,
                "source": "RAG",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "RAG 查詢未找到相關文檔",
                "error": "no_relevant_documents"
            }

    except Exception as e:
        # 將 RAG 失敗從錯誤級別降級為警告級別
        logger.warning(f"RAG 查詢失敗，但不影響報告生成: {str(e)}")

        # 檢查錯誤類型
        error_type = "query_failed"
        if "shapes" in str(e) and "not aligned" in str(e):
            error_type = "vector_dimension_mismatch"
            logger.warning("檢測到向量維度不匹配錯誤，這通常是 RAG 服務配置問題")
        elif "connection" in str(e).lower() or "timeout" in str(e).lower():
            error_type = "connection_failed"
            logger.warning("RAG 服務連接失敗，可能服務不可用")

        return {
            "ok": False,
            "source": "RAG",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"RAG 查詢不可用: {str(e)}",
            "error": error_type,
            "fallback_note": "RAG 資料源暫時不可用，報告將基於其他可用資料源生成"
        }


# ===== LLM 驅動的智能查詢解析 =====

async def parse_report_query_with_llm(query: str, llm=None) -> Dict[str, Any]:
    """使用 LLM 智能解析報告查詢"""
    logger.info(f"使用 LLM 解析報告查詢: {query}")

    if not llm:
        logger.warning("LLM 未設定，回退到規則式解析")
        return parse_report_query_fallback(query)

    try:
        # 構建結構化提示
        system_prompt = """你是專業的金融查詢分析師。請分析用戶的報告請求，並返回結構化的 JSON 響應。

請根據用戶查詢內容，智能判斷以下信息：
1. report_type: 報告類型 ("stock", "macro", "news", "custom")
2. symbols: 股票代號列表 (如 ["AAPL", "TSLA"])
3. indicators: 總經指標列表 (如 ["GDP", "CPI", "UNEMPLOYMENT"])
4. keywords: 關鍵詞列表
5. confidence: 判斷信心度 (0.0-1.0)

判斷規則：
- 如果提到股票代號、公司名稱、股票分析 → stock
- 如果提到 GDP、CPI、通膨、失業率、總經 → macro
- 如果提到新聞、消息、資訊 → news
- 其他情況 → custom

請只返回 JSON 格式，不要其他文字。"""

        user_prompt = f"請分析這個報告請求：{query}"

        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=user_prompt)
        ]

        # 調用 LLM
        response = await llm.ainvoke(messages, temperature=0.1, max_tokens=500)

        # 解析 JSON 響應
        try:
            result = json.loads(response.content.strip())

            # 驗證必要欄位
            if not isinstance(result, dict):
                raise ValueError("響應不是有效的字典格式")

            # 設定預設值
            parsed_result = {
                "report_type": result.get("report_type", "custom"),
                "symbols": result.get("symbols", []),
                "indicators": result.get("indicators", []),
                "keywords": result.get("keywords", []),
                "confidence": result.get("confidence", 0.8),
                "parsing_method": "llm"
            }

            logger.info(f"LLM 解析結果: {parsed_result}")
            return parsed_result

        except (json.JSONDecodeError, ValueError) as e:
            logger.warning(f"LLM 響應解析失敗: {e}, 響應內容: {response.content}")
            return parse_report_query_fallback(query)

    except Exception as e:
        logger.error(f"LLM 查詢解析失敗: {str(e)}")
        return parse_report_query_fallback(query)


def parse_report_query_fallback(query: str) -> Dict[str, Any]:
    """規則式查詢解析 (備用方案)"""
    logger.info(f"使用規則式解析報告查詢: {query}")

    # 移除 /report 前綴
    content = query.replace("/report", "").strip()

    # 預設值
    result = {
        "report_type": "custom",
        "symbols": [],
        "indicators": [],
        "keywords": [],
        "confidence": 0.6,
        "parsing_method": "fallback"
    }

    # 檢測報告類型
    if re.search(r'\b(stock|stocks|股票|個股|分析)\b', content, re.IGNORECASE):
        result["report_type"] = "stock"
    elif re.search(r'\b(macro|總經|經濟|gdp|cpi|inflation)\b', content, re.IGNORECASE):
        result["report_type"] = "macro"
    elif re.search(r'\b(news|新聞|消息)\b', content, re.IGNORECASE):
        result["report_type"] = "news"

    # 提取股票代號
    symbols = re.findall(r'\b[A-Z]{2,5}\b', content.upper())
    exclude_words = {'STOCK', 'STOCKS', 'MACRO', 'NEWS', 'CUSTOM', 'REPORT', 'GDP', 'CPI'}
    valid_symbols = [s for s in symbols if s not in exclude_words]

    if valid_symbols and result["report_type"] == "custom":
        result["report_type"] = "stock"
        logger.info(f"檢測到股票代號 {valid_symbols}，自動設定為股票報告")

    result["symbols"] = list(set(valid_symbols))

    # 提取總經指標關鍵詞
    macro_keywords = re.findall(r'\b(GDP|CPI|INFLATION|UNEMPLOYMENT|INTEREST|RATE)\b', content.upper())
    result["indicators"] = list(set(macro_keywords))

    # 其他關鍵詞
    result["keywords"] = content.split()

    logger.info(f"規則式解析結果: {result}")
    return result


# ===== 報告建構工具 =====

async def tool_build_report(template_id: str, context: Dict[str, Any], output_formats: List[str] = None) -> Dict[str, Any]:
    """報告建構工具"""
    _ensure_imports()
    logger.info(f"建構報告: {template_id}, 格式: {output_formats}")

    if output_formats is None:
        output_formats = ["markdown", "pdf"]

    try:
        from app.services.report import report_service
        from app.settings import settings
        import markdown
        import weasyprint
        from pathlib import Path

        # 生成時間戳和檔案名
        timestamp = generate_timestamp()
        slug = context.get("slug", "report")

        # 設定輸出目錄
        output_dir = Path(settings.output_dir) / "reports"
        output_dir.mkdir(parents=True, exist_ok=True)

        # 準備模板變數
        template_vars = {
            "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "timestamp": timestamp,
            "slug": slug,
            **context
        }

        # 載入模板
        try:
            template = report_service.get_template(template_id)
            if not template:
                return {
                    "ok": False,
                    "source": "REPORT",
                    "data": None,
                    "timestamp": datetime.now().isoformat(),
                    "logs": f"找不到模板: {template_id}",
                    "error": "template_not_found"
                }
        except Exception as e:
            return {
                "ok": False,
                "source": "REPORT",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"載入模板失敗: {str(e)}",
                "error": "template_load_failed"
            }

        # 渲染模板
        try:
            rendered_content = template.render(**template_vars)
        except Exception as e:
            return {
                "ok": False,
                "source": "REPORT",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"模板渲染失敗: {str(e)}",
                "error": "template_render_failed"
            }

        output_files = []

        # 生成 Markdown 檔案
        if "markdown" in output_formats:
            md_file = output_dir / f"{slug}.md"
            try:
                with open(md_file, 'w', encoding='utf-8') as f:
                    f.write(rendered_content)

                output_files.append({
                    "format": "markdown",
                    "filename": f"{slug}.md",
                    "path": str(md_file),
                    "size": md_file.stat().st_size
                })
                logger.info(f"已生成 Markdown 檔案: {md_file}")
            except Exception as e:
                logger.error(f"生成 Markdown 檔案失敗: {str(e)}")

        # 生成 PDF 檔案
        if "pdf" in output_formats:
            try:
                pdf_file = output_dir / f"{slug}.pdf"

                # 轉換 Markdown 為 HTML
                html_content = markdown.markdown(rendered_content)

                # 載入 CSS 樣式
                css_path = Path(settings.pdf_default_css)
                css_content = ""
                if css_path.exists():
                    with open(css_path, 'r', encoding='utf-8') as f:
                        css_content = f.read()

                # 完整 HTML
                full_html = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{slug} 報告</title>
    <style>{css_content}</style>
</head>
<body>
    {html_content}
</body>
</html>
"""

                # 生成 PDF
                weasyprint.HTML(string=full_html).write_pdf(str(pdf_file))

                output_files.append({
                    "format": "pdf",
                    "filename": f"{slug}.pdf",
                    "path": str(pdf_file),
                    "size": pdf_file.stat().st_size
                })
                logger.info(f"已生成 PDF 檔案: {pdf_file}")

            except ImportError:
                logger.warning("weasyprint 未安裝，跳過 PDF 生成")
            except Exception as e:
                logger.error(f"生成 PDF 檔案失敗: {str(e)}")

        # 生成 PPTX 檔案（可選）
        if "pptx" in output_formats:
            try:
                from pptx import Presentation

                pptx_file = output_dir / f"{slug}.pptx"

                # 建立簡單的 PowerPoint
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = f"{slug} 報告"
                content.text = f"生成時間: {template_vars['generated_at']}\n\n{rendered_content[:500]}..."

                prs.save(str(pptx_file))

                output_files.append({
                    "format": "pptx",
                    "filename": f"{slug}.pptx",
                    "path": str(pptx_file),
                    "size": pptx_file.stat().st_size
                })
                logger.info(f"已生成 PPTX 檔案: {pptx_file}")

            except ImportError:
                logger.warning("python-pptx 未安裝，跳過 PPTX 生成")
            except Exception as e:
                logger.error(f"生成 PPTX 檔案失敗: {str(e)}")

        return {
            "ok": True,
            "source": "REPORT",
            "data": {
                "files": output_files,
                "template_id": template_id,
                "slug": slug,
                "content_preview": rendered_content[:200] + "..." if len(rendered_content) > 200 else rendered_content
            },
            "timestamp": datetime.now().isoformat(),
            "logs": f"成功生成 {len(output_files)} 個檔案"
        }

    except Exception as e:
        logger.error(f"報告建構失敗: {str(e)}")
        return {
            "ok": False,
            "source": "REPORT",
            "data": None,
            "timestamp": datetime.now().isoformat(),
            "logs": f"報告建構失敗: {str(e)}",
            "error": "build_failed"
        }


# ===== 輔助函數 =====

def generate_timestamp() -> str:
    """生成時間戳記"""
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def generate_slug(symbols: List[str], report_type: str) -> str:
    """生成檔案名稱 slug"""
    if symbols:
        # 取前兩個 symbol
        slug_symbols = symbols[:2]
        return "_".join(slug_symbols)
    else:
        return "CUSTOM" if report_type == "custom" else report_type.upper()


# ===== 簡化報告代理類 =====

class SimpleReportAgent:
    """簡化報告代理 - 使用線性處理流程和 LLM 自主決策"""

    def __init__(self):
        _ensure_imports()
        self.llm = self._create_llm()
        logger.info("簡化報告代理初始化完成")

    def _create_llm(self):
        """建立 LLM 實例"""
        try:
            from app.settings import settings
            from langchain_openai import ChatOpenAI

            if settings.openai_api_key:
                return ChatOpenAI(
                    model="gpt-4o-mini",
                    api_key=settings.openai_api_key,
                    temperature=0.3
                )
            elif settings.azure_openai_api_key and settings.azure_openai_endpoint:
                return ChatOpenAI(
                    model=settings.azure_openai_deployment,
                    api_key=settings.azure_openai_api_key,
                    azure_endpoint=settings.azure_openai_endpoint,
                    api_version=settings.azure_openai_api_version,
                    temperature=0.3
                )
            else:
                logger.warning("未設定 OpenAI 或 Azure OpenAI API 金鑰，LLM 增強功能將無法使用")
                return None
        except ImportError:
            logger.warning("未安裝 langchain_openai，LLM 增強功能將無法使用")
            return None

    async def run(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """執行報告生成 - 線性處理流程"""
        try:
            logger.info("開始執行簡化報告生成流程")

            # 1. 解析查詢
            query = input_data["query"]
            session_id = input_data.get("session_id", "simple_report")

            # 使用 LLM 智能解析查詢
            parsed_query = await parse_report_query_with_llm(query, self.llm)

            # 生成時間戳和 slug
            timestamp = generate_timestamp()
            slug = generate_slug(parsed_query["symbols"], parsed_query["report_type"])

            logger.info(f"查詢解析完成: 類型={parsed_query['report_type']}, 股票代號={parsed_query['symbols']}")

            # 2. 收集資料
            context = await self._collect_data(parsed_query)
            context["parsed_query"] = parsed_query
            context["slug"] = slug
            context["symbols"] = parsed_query["symbols"]

            logger.info("資料收集完成")

            # 3. 選擇模板
            template_id = self._select_template(parsed_query["report_type"])

            logger.info(f"已選擇模板: {template_id}")

            # 4. LLM 增強分析（如果啟用）
            if self.llm:
                try:
                    enhanced_context = await self._enhance_context_with_llm(context, template_id)
                    context.update(enhanced_context)
                    logger.info("LLM 增強分析完成")
                except Exception as e:
                    logger.warning(f"LLM 增強分析失敗，使用原始資料: {str(e)}")

            # 5. 建構報告
            output_formats = ["markdown", "pdf"]
            report_result = await tool_build_report(template_id, context, output_formats)

            if not report_result.get("ok"):
                raise Exception(f"報告建構失敗: {report_result.get('logs', '未知錯誤')}")

            # 6. 建構回應
            output_files = report_result["data"]["files"]

            # 生成回應訊息
            response_message = self._generate_response_message(parsed_query, output_files, context)

            return {
                "ok": True,
                "response": response_message,
                "input_type": input_data["input_type"],
                "query": query,
                "output_files": output_files,
                "session_id": session_id,
                "timestamp": datetime.now().isoformat(),
                "trace_id": input_data.get("trace_id", "")
            }

        except Exception as e:
            logger.error(f"簡化報告代理執行失敗: {str(e)}")
            return {
                "ok": False,
                "response": f"報告生成失敗：{str(e)}",
                "input_type": input_data["input_type"],
                "query": input_data["query"],
                "output_files": [],
                "session_id": input_data.get("session_id"),
                "timestamp": datetime.now().isoformat(),
                "trace_id": input_data.get("trace_id", ""),
                "error": str(e)
            }

    async def _collect_data(self, parsed_query: Dict[str, Any]) -> Dict[str, Any]:
        """收集資料"""
        from app.settings import settings

        context = {}
        report_type = parsed_query["report_type"]
        symbols = parsed_query["symbols"]

        # 檢查是否啟用工具執行
        if not settings.execute_tools:
            logger.warning("工具執行已停用")
            context["error"] = {
                "ok": False,
                "source": "SYSTEM",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": "工具執行已停用",
                "error": "execute_tools_disabled"
            }
            return context

        try:
            # 根據報告類型收集資料
            if report_type == "stock" and symbols:
                logger.info(f"收集股票資料: {symbols}")

                # 獲取報價
                quotes_result = await _agent_tool_fmp_quote.ainvoke({"symbols": symbols})
                context["quotes"] = self._wrap_tool_result(quotes_result, "FMP")

                # 獲取公司基本面
                profiles_result = await _agent_tool_fmp_profile.ainvoke({"symbols": symbols})
                context["profiles"] = self._wrap_tool_result(profiles_result, "FMP")

                # 獲取新聞
                news_result = await _agent_tool_fmp_news.ainvoke({"symbols": symbols, "limit": 10})
                context["news"] = self._wrap_tool_result(news_result, "FMP")

            elif report_type == "macro":
                logger.info("收集總經資料")
                indicators = parsed_query.get("indicators", ["GDP", "CPI", "UNEMPLOYMENT"])

                # 收集總經資料
                macro_data = {}
                for indicator in indicators:
                    try:
                        result = await _agent_tool_fmp_macro.ainvoke({"indicator": indicator, "country": "US"})
                        if result.get("ok"):
                            macro_data[f"{indicator}_US"] = result.get("data", [])
                        else:
                            logger.warning(f"總經指標 {indicator} 查詢失敗")
                            macro_data[f"{indicator}_US"] = []
                    except Exception as e:
                        logger.error(f"總經指標 {indicator} 查詢異常: {str(e)}")
                        macro_data[f"{indicator}_US"] = []

                context["macro"] = {
                    "ok": True,
                    "source": "FMP",
                    "data": macro_data,
                    "timestamp": datetime.now().isoformat(),
                    "logs": f"查詢 {len(indicators)} 個總經指標"
                }

                # 也收集相關新聞
                news_result = await _agent_tool_fmp_news.ainvoke({"symbols": [], "limit": 5})
                context["news"] = self._wrap_tool_result(news_result, "FMP")

            elif report_type == "news":
                logger.info("收集新聞資料")
                news_result = await _agent_tool_fmp_news.ainvoke({"symbols": symbols, "limit": 15})
                context["news"] = self._wrap_tool_result(news_result, "FMP")

            else:
                logger.info("自訂報告，收集基本資料")
                if symbols:
                    quotes_result = await _agent_tool_fmp_quote.ainvoke({"symbols": symbols})
                    context["quotes"] = self._wrap_tool_result(quotes_result, "FMP")

            # 嘗試 RAG 查詢（可選）
            try:
                keywords = " ".join(parsed_query.get("keywords", []))
                if keywords:
                    rag_result = await _agent_tool_rag_query.ainvoke({"question": keywords, "top_k": 5})
                    context["rag"] = self._wrap_tool_result(rag_result, "RAG")
            except Exception as e:
                logger.warning(f"RAG 查詢失敗: {str(e)}")
                context["rag"] = {
                    "ok": False,
                    "source": "RAG",
                    "data": None,
                    "timestamp": datetime.now().isoformat(),
                    "logs": f"RAG 查詢失敗: {str(e)}",
                    "error": "query_failed"
                }

        except Exception as e:
            logger.warning(f"資料收集過程中發生異常: {str(e)}")
            context["error"] = {
                "ok": False,
                "source": "SYSTEM",
                "data": None,
                "timestamp": datetime.now().isoformat(),
                "logs": f"資料收集異常: {str(e)}",
                "error": "data_collection_failed"
            }

        return context

    def _wrap_tool_result(self, result: Dict[str, Any], source: str) -> Dict[str, Any]:
        """包裝工具結果以保持格式一致性"""
        return {
            "ok": result.get("ok", False),
            "source": source,
            "data": result.get("data"),
            "timestamp": datetime.now().isoformat(),
            "logs": result.get("logs", ""),
            "error": result.get("reason") if not result.get("ok") else None
        }

    def _select_template(self, report_type: str) -> str:
        """選擇報告模板"""
        from app.settings import settings

        # 模板映射
        template_mapping = {
            "stock": "stock.j2",
            "macro": "macro.j2",
            "news": "news.j2",
            "custom": "custom.j2"
        }

        base_template_id = template_mapping.get(report_type, "custom.j2")

        # 檢查是否啟用 LLM 增強並選擇對應模板
        if settings.llm_report_enhancement and report_type == "stock":
            # 對於股票報告，如果啟用 LLM 增強，使用增強模板
            template_id = "stock_llm_enhanced.j2"
            logger.info(f"LLM 增強已啟用，使用增強模板: {template_id}")
        elif report_type == "stock":
            # 即使 LLM 增強未啟用，股票報告也使用增強模板（提供更好的結構）
            template_id = "stock_llm_enhanced.j2"
            logger.info(f"股票報告使用增強模板（結構更完整）: {template_id}")
        else:
            template_id = base_template_id
            logger.info(f"使用基礎模板: {template_id}")

        return template_id

    async def _enhance_context_with_llm(self, context: Dict[str, Any], template_id: str) -> Dict[str, Any]:
        """使用 LLM 增強報告上下文 - 基於實際資料進行深度分析"""
        try:
            if not self.llm:
                logger.warning("LLM 未設定，跳過增強分析")
                return {}

            # 檢查是否有足夠的資料進行分析
            has_data = any([
                context.get("quotes", {}).get("ok"),
                context.get("profiles", {}).get("ok"),
                context.get("news", {}).get("ok"),
                context.get("macro", {}).get("ok")
            ])

            if not has_data:
                logger.warning("沒有可用的資料進行 LLM 分析")
                return {
                    "market_analysis": "由於缺乏資料，無法進行詳細的市場分析",
                    "fundamental_analysis": "需要更多公司基本面資料才能進行分析",
                    "news_impact": "沒有相關新聞資料可供分析",
                    "investment_recommendation": "資料不足，建議謹慎投資",
                    "risk_assessment": "由於資料限制，風險評估可能不完整",
                    "key_insights": ["資料收集受限，建議獲取更多資訊後再做決策"]
                }

            # 建構詳細的分析提示
            prompt = self._build_analysis_prompt(context, template_id)
            logger.info(f"建構 LLM 分析提示，長度: {len(prompt)} 字符")

            # 建立訊息
            system_prompt = """你是資深的金融分析師，具有豐富的股票分析經驗。

請基於提供的實際資料進行專業分析：
1. 僅使用提供的真實資料，不要杜撰任何數據
2. 提供具體、可操作的分析和建議
3. 引用具體的數據點支持你的分析
4. 保持客觀和專業的分析態度
5. 必須返回有效的 JSON 格式

分析應該具體且實用，避免使用模糊或通用的表述。"""

            messages = [
                SystemMessage(content=system_prompt),
                HumanMessage(content=prompt)
            ]

            # 調用 LLM
            from app.settings import settings
            logger.info("開始調用 LLM 進行深度分析...")

            response = await self.llm.ainvoke(
                messages,
                temperature=settings.llm_analysis_temperature,
                max_tokens=settings.llm_analysis_max_tokens
            )

            logger.info(f"LLM 響應長度: {len(response.content)} 字符")

            # 解析 LLM 響應
            try:
                # 清理響應內容，移除可能的 markdown 標記
                content = response.content.strip()
                if content.startswith("```json"):
                    content = content[7:]
                if content.endswith("```"):
                    content = content[:-3]
                content = content.strip()

                analysis_result = json.loads(content)
                logger.info("成功解析 LLM JSON 響應")

                # 驗證必要欄位
                required_fields = ["market_analysis", "fundamental_analysis", "news_impact",
                                 "investment_recommendation", "risk_assessment", "key_insights"]

                for field in required_fields:
                    if field not in analysis_result:
                        logger.warning(f"LLM 響應缺少必要欄位: {field}")
                        analysis_result[field] = f"分析中缺少 {field} 內容"

                # 提取增強資料
                enhanced_data = {
                    "llm_analysis": analysis_result,
                    "market_analysis": analysis_result.get("market_analysis", ""),
                    "fundamental_analysis": analysis_result.get("fundamental_analysis", ""),
                    "news_impact": analysis_result.get("news_impact", ""),
                    "investment_recommendation": analysis_result.get("investment_recommendation", ""),
                    "risk_assessment": analysis_result.get("risk_assessment", ""),
                    "key_insights": analysis_result.get("key_insights", [])
                }

                logger.info("LLM 增強分析成功完成")
                return enhanced_data

            except json.JSONDecodeError as e:
                logger.warning(f"LLM 響應不是有效的 JSON: {str(e)}")
                logger.warning(f"原始響應內容: {response.content[:500]}...")

                # 嘗試從文字中提取有用信息
                content = response.content
                return {
                    "llm_analysis": content,
                    "market_analysis": f"基於 LLM 分析：{content[:200]}..." if len(content) > 200 else content,
                    "fundamental_analysis": "LLM 分析結果格式異常，請檢查原始分析內容",
                    "news_impact": "無法解析新聞影響評估",
                    "investment_recommendation": "建議查看原始 LLM 分析內容",
                    "risk_assessment": "風險評估解析失敗",
                    "key_insights": ["LLM 響應格式異常", "建議檢查原始分析內容"]
                }

        except Exception as e:
            logger.error(f"LLM 增強分析失敗: {str(e)}")
            import traceback
            logger.error(f"錯誤詳情: {traceback.format_exc()}")

            return {
                "market_analysis": f"LLM 分析過程中發生錯誤: {str(e)}",
                "fundamental_analysis": "分析服務暫時不可用",
                "news_impact": "無法評估新聞影響",
                "investment_recommendation": "建議等待分析服務恢復後再做決策",
                "risk_assessment": "由於技術問題，風險評估不可用",
                "key_insights": ["分析服務異常", "建議稍後重試"]
            }

    def _build_analysis_prompt(self, context: Dict[str, Any], template_id: str) -> str:
        """建構 LLM 分析提示 - 包含實際資料內容"""
        symbols = context.get("symbols", [])

        prompt_parts = [
            f"你是專業的金融分析師。請基於以下實際資料進行深度分析，並提供結構化的 JSON 響應：",
            f"股票代號: {', '.join(symbols) if symbols else '無'}",
            f"模板類型: {template_id}",
            f"分析時間: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "",
            "=== 實際資料內容 ==="
        ]

        # 添加股價資料詳情
        if context.get("quotes", {}).get("ok") and context["quotes"].get("data"):
            prompt_parts.append("\n📊 股價資料:")
            quotes_data = context["quotes"]["data"]
            if isinstance(quotes_data, list):
                for quote in quotes_data[:3]:  # 限制前3個
                    if isinstance(quote, dict):
                        symbol = quote.get("symbol", "N/A")
                        price = quote.get("price", "N/A")
                        change = quote.get("change", "N/A")
                        change_pct = quote.get("changesPercentage", "N/A")
                        volume = quote.get("volume", "N/A")
                        day_low = quote.get("dayLow", "N/A")
                        day_high = quote.get("dayHigh", "N/A")

                        prompt_parts.append(f"- {symbol}: 價格 ${price}, 漲跌 {change} ({change_pct}%), 成交量 {volume}")
                        prompt_parts.append(f"  日內區間: ${day_low} - ${day_high}")

        # 添加公司基本面資料詳情
        if context.get("profiles", {}).get("ok") and context["profiles"].get("data"):
            prompt_parts.append("\n🏢 公司基本面:")
            profiles_data = context["profiles"]["data"]
            if isinstance(profiles_data, list):
                for profile in profiles_data[:3]:  # 限制前3個
                    if isinstance(profile, dict):
                        symbol = profile.get("symbol", "N/A")
                        company_name = profile.get("companyName", "N/A")
                        industry = profile.get("industry", "N/A")
                        sector = profile.get("sector", "N/A")
                        market_cap = profile.get("mktCap", "N/A")
                        description = profile.get("description", "")[:200] + "..." if profile.get("description") else "N/A"

                        prompt_parts.append(f"- {symbol} ({company_name})")
                        prompt_parts.append(f"  產業: {industry}, 部門: {sector}")
                        prompt_parts.append(f"  市值: {market_cap}")
                        prompt_parts.append(f"  描述: {description}")

        # 添加新聞資料詳情
        if context.get("news", {}).get("ok") and context["news"].get("data"):
            prompt_parts.append("\n📰 相關新聞:")
            news_data = context["news"]["data"]
            if isinstance(news_data, list):
                for i, news in enumerate(news_data[:5], 1):  # 限制前5條新聞
                    if isinstance(news, dict):
                        title = news.get("title", "N/A")
                        text = news.get("text", "")[:150] + "..." if news.get("text") else "N/A"
                        published_date = news.get("publishedDate", "N/A")
                        site = news.get("site", "N/A")

                        prompt_parts.append(f"{i}. {title}")
                        prompt_parts.append(f"   來源: {site}, 時間: {published_date}")
                        prompt_parts.append(f"   摘要: {text}")

        # 添加總經資料詳情（如果有）
        if context.get("macro", {}).get("ok") and context["macro"].get("data"):
            prompt_parts.append("\n📈 總經資料:")
            macro_data = context["macro"]["data"]
            if isinstance(macro_data, dict):
                for indicator, data in macro_data.items():
                    if data and isinstance(data, list):
                        latest = data[0] if data else {}
                        if isinstance(latest, dict):
                            value = latest.get("value", "N/A")
                            date = latest.get("date", "N/A")
                            prompt_parts.append(f"- {indicator}: {value} (截至 {date})")

        prompt_parts.extend([
            "",
            "=== 分析要求 ===",
            "請基於上述實際資料，提供深度且具體的分析。避免使用通用或模糊的表述。",
            "每個分析項目應該：",
            "1. 引用具體的數據點",
            "2. 提供明確的判斷和理由",
            "3. 給出可操作的建議",
            "",
            "請提供以下 JSON 格式的分析：",
            "{",
            '  "market_analysis": "基於股價表現、成交量、漲跌幅等數據的具體市場分析",',
            '  "fundamental_analysis": "基於公司基本面資料的具體分析，包括產業地位、市值評估等",',
            '  "news_impact": "基於實際新聞內容的影響評估，分析正面/負面因素",',
            '  "investment_recommendation": "基於綜合分析的具體投資建議，包括目標價位或操作策略",',
            '  "risk_assessment": "基於當前市況和公司狀況的具體風險評估",',
            '  "key_insights": ["基於數據的關鍵洞察1", "基於數據的關鍵洞察2", "基於數據的關鍵洞察3"]',
            "}"
        ])

        return "\n".join(prompt_parts)

    def _generate_response_message(self, parsed_query: Dict[str, Any], output_files: List[Dict[str, Any]], context: Dict[str, Any]) -> str:
        """生成回應訊息"""
        report_type = parsed_query["report_type"]
        symbols = parsed_query["symbols"]

        # 基本訊息
        message_parts = [
            f"📊 {report_type.upper()} 報告生成完成！",
            ""
        ]

        # 報告詳情
        if symbols:
            message_parts.append(f"🎯 分析標的：{', '.join(symbols)}")

        message_parts.append(f"📝 報告類型：{report_type}")
        message_parts.append(f"🕒 生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        message_parts.append("")

        # 資料來源摘要
        data_sources = []
        for key, value in context.items():
            if isinstance(value, dict) and "ok" in value:
                status = "✅" if value["ok"] else "❌"
                source_name = value.get("source", key.upper())
                data_sources.append(f"{status} {source_name}")

        if data_sources:
            message_parts.append("📊 資料來源：")
            message_parts.extend([f"  {source}" for source in data_sources])
            message_parts.append("")

        # 輸出檔案
        if output_files:
            message_parts.append("📄 生成檔案：")
            for file_info in output_files:
                format_name = file_info["format"].upper()
                file_path = file_info["path"]
                message_parts.append(f"  • {format_name}: {file_path}")
            message_parts.append("")

        # 解析方法資訊
        parsing_method = parsed_query.get("parsing_method", "unknown")
        confidence = parsed_query.get("confidence", 0.0)

        if parsing_method == "llm":
            message_parts.append(f"🤖 智能解析：LLM 驅動 (信心度: {confidence:.1%})")
        else:
            message_parts.append(f"🔧 規則解析：備用方案 (信心度: {confidence:.1%})")

        message_parts.append("")
        message_parts.append("✨ 報告已準備就緒，請查看生成的檔案！")

        return "\n".join(message_parts)


# ===== 全域實例 =====

# 全域簡化報告代理實例（延遲初始化）
simple_report_agent = None


def get_simple_report_agent() -> SimpleReportAgent:
    """取得全域簡化報告代理實例"""
    global simple_report_agent
    if simple_report_agent is None:
        simple_report_agent = SimpleReportAgent()
    return simple_report_agent


# ===== 向後兼容性支援 =====

class ReportAgent:
    """向後兼容的報告代理類 - 委託給 SimpleReportAgent"""

    def __init__(self):
        self.simple_agent = SimpleReportAgent()
        logger.info("向後兼容報告代理初始化完成")

    async def run(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """執行報告生成 - 委託給簡化代理"""
        return await self.simple_agent.run(input_data)


# 全域報告代理實例（向後兼容）
report_agent = None


def get_report_agent() -> ReportAgent:
    """取得全域報告代理實例（向後兼容）"""
    global report_agent
    if report_agent is None:
        report_agent = ReportAgent()
    return report_agent


# ===== 直接執行支援 =====

async def run_default_report_generation():
    """執行預設報告生成（無需命令列參數）"""
    global simple_report_agent

    # 初始化簡化報告代理（如果尚未初始化）
    if simple_report_agent is None:
        simple_report_agent = SimpleReportAgent()

    # 硬編碼的預設值
    DEFAULT_QUERY = "/report stock AAPL TSLA NVDA"
    DEFAULT_OUTPUT_DIR = "./outputs"
    DEFAULT_FORMATS = ["markdown", "pdf"]

    print("🤖 簡化報告代理 - 預設報告生成")
    print(f"📝 查詢：{DEFAULT_QUERY}")
    print(f"📁 輸出目錄：{DEFAULT_OUTPUT_DIR}")
    print(f"📄 輸出格式：{', '.join(DEFAULT_FORMATS)}")
    print()

    try:
        # 準備輸入資料
        input_data = {
            "input_type": "text",
            "query": DEFAULT_QUERY,
            "session_id": f"simple_report_{generate_timestamp()}",
            "trace_id": f"simple_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        }

        print("🚀 開始執行報告生成...")

        # 執行報告生成
        result = await simple_report_agent.run(input_data)

        # 處理結果
        if result.get("ok"):
            print("✅ 報告生成成功！")
            print()

            # 顯示輸出檔案
            output_files = result.get("output_files", [])
            if output_files:
                print("📄 生成的檔案：")
                for file_info in output_files:
                    file_path = file_info.get("path", "未知路徑")
                    file_format = file_info.get("format", "未知格式")
                    print(f"  • {file_format.upper()}: {file_path}")

            # 顯示回應內容
            if result.get("response"):
                print()
                print("📋 報告摘要：")
                print(result["response"])

            return 0

        else:
            print("❌ 報告生成失敗")
            error_msg = result.get("error", "未知錯誤")
            print(f"錯誤：{error_msg}")

            if result.get("response"):
                print()
                print("詳細資訊：")
                print(result["response"])

            return 1

    except KeyboardInterrupt:
        print("\n⏹️ 使用者中斷執行")
        return 130

    except Exception as e:
        print(f"❌ 執行過程中發生錯誤：{str(e)}")
        return 1


# ===== 主程式入口 =====

if __name__ == "__main__":
    """直接執行支援"""
    import asyncio

    async def main():
        """主程式"""
        try:
            exit_code = await run_default_report_generation()
            return exit_code
        except Exception as e:
            print(f"❌ 程式執行失敗：{str(e)}")
            return 1

    # 執行主程式
    exit_code = asyncio.run(main())
    exit(exit_code)